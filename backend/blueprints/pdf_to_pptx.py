from io import BytesIO
import fitz
from pptx import Presentation
from pptx.util import Inches, Emu
from flask import Blueprint, request

from utils.helpers import error, send_file_and_cleanup
from utils.validators import validate_uploaded_file, validate_pdf_file

pdf_pptx_bp = Blueprint("pdf_pptx", __name__)


@pdf_pptx_bp.route("/convertPdfToPptx", methods=["POST"])
def convert_pdf_to_pptx():
    doc = None
    pdf_bytes = None

    try:
        pdf_file, filename, upload_error = validate_uploaded_file(
            request,
            "file",
        )

        if upload_error:
            return upload_error

        pdf_error = validate_pdf_file(pdf_file, filename)

        if pdf_error:
            return pdf_error

        pdf_bytes = pdf_file.read()

        doc = fitz.open(
            stream=pdf_bytes,
            filetype="pdf",
        )

        try:
            if doc.page_count == 0:
                return error("Empty PDF")

            prs = Presentation()

            for page_num in range(doc.page_count):
                page = doc.load_page(page_num)

                page_width = page.rect.width
                page_height = page.rect.height

                aspect_ratio = page_width / page_height

                max_width_inches = 13.333
                max_height_inches = 7.5

                if aspect_ratio >= (max_width_inches / max_height_inches):
                    slide_width = Inches(max_width_inches)
                    slide_height = Inches(max_width_inches / aspect_ratio)
                else:
                    slide_height = Inches(max_height_inches)
                    slide_width = Inches(max_height_inches * aspect_ratio)

                prs.slide_width = slide_width
                prs.slide_height = slide_height

                zoom = max_width_inches * 96 / page_width
                mat = fitz.Matrix(zoom, zoom)

                pix = page.get_pixmap(
                    matrix=mat,
                    alpha=False,
                )

                img_bytes = pix.tobytes(output="png")

                img_stream = BytesIO(img_bytes)

                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)

                slide.shapes.add_picture(
                    img_stream,
                    Inches(0),
                    Inches(0),
                    width=slide_width,
                    height=slide_height,
                )

                pix = None

            output_stream = BytesIO()
            prs.save(output_stream)
            output_stream.seek(0)

            import gc
            gc.collect()

            return send_file_and_cleanup(
                output_stream,
                mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                as_attachment=True,
                download_name="presentation.pptx",
            )

        finally:
            if doc:
                doc.close()

    except Exception as e:
        import traceback
        traceback.print_exc()
        return error(
            f"Conversion failed: {str(e)}",
            500,
        )